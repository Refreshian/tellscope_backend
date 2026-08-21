from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import insights
from .catalog import Catalog
from .metrics import cooccurrence_media, daily_series
from .models import ReportBundle

NAVY = RGBColor(0x1A, 0x36, 0x5D)
NAVY2 = RGBColor(0x2C, 0x52, 0x82)
RED = RGBColor(0xC4, 0x1E, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5B, 0x67, 0x75)
LINE = RGBColor(0xD6, 0xDC, 0xE2)
CARD = RGBColor(0xF4, 0xF6, 0xF8)
BLUE = RGBColor(0x44, 0x72, 0xC4)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x70, 0xAD, 0x47)

W = Inches(13.333)
H = Inches(7.5)


def _bullets(settings, key, fallback):
    ov = (settings or {}).get("insight_overrides") or {}
    val = ov.get(key)
    if isinstance(val, list) and val:
        return [str(x) for x in val[:4]]
    return fallback


def _set_run(run, text: str, size=14, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _textbox(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    return box


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, color)
    return sh


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, WHITE)
    _rect(slide, 0, 0, W, Inches(0.08), RED)
    return slide


def _header(slide, kicker: str, title: str, page: str, period: str):
    _rect(slide, 0, Inches(0.08), W, Inches(0.78), NAVY)
    _textbox(slide, Inches(0.4), Inches(0.12), Inches(10), Inches(0.28), kicker.upper(), 11, True, RGBColor(0xF2, 0xC9, 0x4C))
    _textbox(slide, Inches(0.4), Inches(0.36), Inches(11.2), Inches(0.44), title, 22, True, WHITE)
    _textbox(slide, Inches(11.6), Inches(0.22), Inches(1.4), Inches(0.44), page, 16, True, WHITE, PP_ALIGN.RIGHT)
    _textbox(slide, Inches(0.4), Inches(7.18), Inches(10), Inches(0.24), period, 10, False, MUTED)
    _textbox(slide, Inches(10.4), Inches(7.18), Inches(2.5), Inches(0.24), "Мосинформ.Рейтинг", 10, False, MUTED, PP_ALIGN.RIGHT)


def _cards(slide, bullets: list[str], top=Inches(5.05)):
    n = max(len(bullets), 1)
    gap = Inches(0.12)
    left = Inches(0.4)
    width = (W - Inches(0.8) - gap * (n - 1)) / n
    for i, text in enumerate(bullets):
        x = left + i * (width + gap)
        _rect(slide, x, top, width, Inches(1.95), CARD)
        _rect(slide, x, top, Inches(0.08), Inches(1.95), RED if i == 0 else NAVY2)
        box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(0.12), width - Inches(0.28), Inches(1.7))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_run(p.add_run(), text, 11, False, INK)


def _bar_chart(slide, left, top, width, height, categories: list[str], values: list[float], title: str = ""):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(title or " ", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    s = plot.series[0]
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = NAVY2
    return chart


def _line_chart(slide, left, top, width, height, categories: list[str], values: list[float]):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Сообщения", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data
    ).chart
    chart.has_legend = False
    return chart


def _metric_slide(prs, kicker, title, page, period, categories, values, bullets, chart_title=""):
    slide = _new_slide(prs)
    _header(slide, kicker, title, page, period)
    if categories and values:
        _bar_chart(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.9), categories, values, chart_title)
    else:
        _textbox(slide, Inches(0.5), Inches(2.2), Inches(12), Inches(1), "Недостаточно данных для графика.", 16, False, MUTED)
    _cards(slide, bullets)
    return slide


def build_pptx(bundle: ReportBundle, catalog: Catalog, path: Path, settings: dict) -> None:
    report = (settings or {}).get("report") or {}
    period = bundle.period_label or ""
    kicker = report.get("kicker") or "Медиаприсутствие ОИВ и городских проектов"
    top_n = int((settings or {}).get("top_n_objects") or 12)
    top = bundle.object_stats[:top_n]
    labels = [catalog.label(s.object_id) for s in top]
    counts = [s.messages for s in top]
    n_msg = len(bundle.messages)
    n_obj = len(bundle.object_stats)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 title
    slide = _new_slide(prs)
    _rect(slide, 0, 0, W, H, NAVY)
    _rect(slide, 0, 0, W, Inches(0.08), RED)
    _textbox(slide, Inches(0.7), Inches(1.8), Inches(12), Inches(0.4), report.get("org_line") or "Мосинформ", 14, False, RGBColor(0xF2, 0xC9, 0x4C))
    _textbox(slide, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2), report.get("product") or "Мосинформ.Рейтинг", 40, True, WHITE)
    _textbox(slide, Inches(0.7), Inches(3.5), Inches(12), Inches(0.6), "Рейтинг медиаприсутствия департаментов и проектов города Москвы", 20, False, WHITE)
    _textbox(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4), period, 18, True, RGBColor(0xF2, 0xC9, 0x4C))
    _textbox(
        slide,
        Inches(0.7),
        Inches(6.4),
        Inches(12),
        Inches(0.4),
        f"{n_msg} уникальных материалов · {n_obj} объектов после разметки",
        14,
        False,
        RGBColor(0xD0, 0xD7, 0xE2),
    )

    # 2 coverage
    slide = _new_slide(prs)
    _header(slide, "Пилот / замер", "Материалы анализа и покрытие выгрузки", "02", period)
    rows = [["Объект", "Сообщений", "Главная роль", "СМИ", "Топ-2", "Охват", "МедиаИндекс"]]
    for st in top[:14]:
        rows.append(
            [
                catalog.label(st.object_id),
                str(st.messages),
                str(st.main_role),
                str(st.unique_media),
                f"{100 * st.top2_share:.0f}%",
                "н/д" if st.reach is None else f"{st.reach / 1_000_000:.1f} млн" if st.reach >= 1_000_000 else str(int(st.reach)),
                "н/д" if st.media_index is None else f"{st.media_index:.0f}",
            ]
        )
    _add_table(slide, rows, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.0))
    _textbox(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.8), " · ".join(bundle.notes[:4] or ["Выгрузка обработана."]), 11, False, MUTED)

    # 3 caveat
    slide = _new_slide(prs)
    _header(slide, "Методологическая оговорка", "Что есть в этой поставке и чего нет", "03", period)
    bullets = [
        f"В отчёт вошло {n_msg} уникальных полных текстов после дедупликации по URL.",
        "Объектная разметка — словарь ОИВ и эвристика; Tellscope подключается для тональности, роли и «инициировано ИЦ».",
        ("Не хватает: " + "; ".join(bundle.missing_metrics) + ".") if bundle.missing_metrics else "Свод Медиалогии подмешан: охват и МедиаИндекс взяты из Excel.",
        "Пилот демонстрирует механику рейтинга. Полноценные позиции — на месячном срезе по согласованному перечню ОИВ.",
    ]
    _cards(slide, bullets, top=Inches(1.3))
    _textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(2.8),
             "Официальный контур — традиционные СМИ и их каналы. Независимый — площадки, где повестку обсуждают без редакционного контура города. Разрыв между контурами в этой версии считается по полю contour у источника.",
             16, False, INK)

    # 4 volume
    _metric_slide(prs, "Параметр 1", "Общий объём упоминаний", "04", period, labels[::-1], counts[::-1], _bullets(settings, "volume", insights.four_volume(bundle, catalog)))

    # 5 speakers
    speaker_labels, speaker_vals = _speaker_bars(bundle, catalog)
    _metric_slide(prs, "Параметр 2", "Упоминание глав и спикеров", "05", period, speaker_labels[::-1], speaker_vals[::-1], _bullets(settings, "speakers", insights.four_speakers(bundle, catalog)))

    # 6 reach
    reach_l = [catalog.label(s.object_id) for s in top if s.reach]
    reach_v = [(s.reach or 0) / 1_000_000 for s in top if s.reach]
    _metric_slide(prs, "Параметр 3", "Суммарный охват", "06", period, reach_l[::-1], reach_v[::-1], _bullets(settings, "reach", insights.four_reach(bundle, catalog)))

    # 7 media table
    slide = _new_slide(prs)
    _header(slide, "Параметр 4", "Распределение по ключевым СМИ", "07", period)
    key = catalog.key_media
    head = ["Объект", *key]
    mrows = [head]
    for st in top[:10]:
        mrows.append([catalog.label(st.object_id), *[str(_media_count(st.media, catalog, name)) for name in key]])
    _add_table(slide, mrows, Inches(0.35), Inches(1.05), Inches(12.6), Inches(3.85))
    _cards(slide, _bullets(settings, "media", insights.four_media(bundle, catalog)), top=Inches(5.05))

    # 8 tone
    pos = [s.positive for s in top]
    neg = [s.negative for s in top]
    neu = [s.neutral for s in top]
    slide = _new_slide(prs)
    _header(slide, "Параметр 5", "Тональность", "08", period)
    data = CategoryChartData()
    data.categories = labels
    data.add_series("Позитив", pos)
    data.add_series("Нейтрал", neu)
    data.add_series("Негатив", neg)
    ch = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.85), data).chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    _cards(slide, _bullets(settings, "tone", insights.four_tone(bundle, catalog)))

    # 9 prominent
    share_main = [100 * s.main_role / max(s.messages, 1) for s in top]
    _metric_slide(prs, "Параметр 6–7", "Доля заметных материалов (главная роль)", "09", period, labels[::-1], share_main[::-1], _bullets(settings, "role", insights.four_role(bundle, catalog)))

    # 10 concentration
    conc = [100 * s.top2_share for s in top]
    _metric_slide(prs, "Параметр 8", "Концентрация по СМИ (доля топ-2)", "10", period, labels[::-1], conc[::-1], _bullets(settings, "concentration", insights.four_concentration(bundle, catalog)))

    # 11 overlap matrix
    slide = _new_slide(prs)
    _header(slide, "Параметр 9–10", "Пересечения по площадкам", "11", period)
    matrix = cooccurrence_media(top, limit=min(8, len(top)))
    names = [catalog.label(s.object_id) for s in top[: len(matrix)]]
    trows = [[""] + names]
    for i, name in enumerate(names):
        trows.append([name] + [str(matrix[i][j]) if i != j else "—" for j in range(len(names))])
    _add_table(slide, trows, Inches(0.4), Inches(1.05), Inches(12.5), Inches(3.85))
    _cards(slide, _bullets(settings, "overlap", insights.four_overlap(bundle, catalog)))

    # 12 initiated
    init_share = [100 * s.initiated / max(s.initiated_known, 1) for s in top]
    _metric_slide(prs, "Параметр 11", "Доля инициированных материалов", "12", period, labels[::-1], init_share[::-1], _bullets(settings, "initiated", insights.four_initiated(bundle, catalog)))

    # 13 dynamics
    slide = _new_slide(prs)
    _header(slide, "Параметр 12", "Динамика публикаций", "13", period)
    series = daily_series(bundle.messages)
    if series:
        cats = [d[5:] for d, _ in series]  # MM-DD
        vals = [n for _, n in series]
        _line_chart(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.9), cats, vals)
        peak_day, peak_n = max(series, key=lambda x: x[1])
        bullets = [
            f"Пик — {peak_day}: {peak_n} материалов.",
            f"Период в текстах: {bundle.period_start:%d.%m} — {bundle.period_end:%d.%m}." if bundle.period_start and bundle.period_end else "",
            "Сравнение дней отражает и динамику повестки, и полноту поставки.",
            f"Всего дней с публикациями: {len(series)}.",
        ]
        _cards(slide, [b for b in bullets if b])
    else:
        _textbox(slide, Inches(0.5), Inches(2), Inches(12), Inches(1), "Нет дат в выгрузке.", 16, False, MUTED)

    # 14 media index
    idx_l = [catalog.label(s.object_id) for s in top if s.media_index]
    idx_v = [s.media_index or 0 for s in top if s.media_index]
    _metric_slide(prs, "Параметр 13", "МедиаИндекс", "14", period, idx_l[::-1], idx_v[::-1], _bullets(settings, "index", insights.four_index(bundle, catalog)))

    # 15 three observations
    slide = _new_slide(prs)
    _header(slide, "Выводы", "Три наблюдения, которые в полной версии становятся решениями", "15", period)
    ov_obs = ((settings or {}).get("insight_overrides") or {}).get("observations")
    obs = ov_obs if isinstance(ov_obs, list) and ov_obs else insights.three_observations(bundle, catalog)
    obs = [
        {
            "title": item.get("title") if isinstance(item, dict) else str(item),
            "fact": item.get("fact", "") if isinstance(item, dict) else "",
            "meaning": item.get("meaning", "") if isinstance(item, dict) else "",
            "full": item.get("full", "") if isinstance(item, dict) else "",
        }
        for item in obs[:3]
    ]
    if not obs:
        obs = insights.three_observations(bundle, catalog)
    for i, item in enumerate(obs):
        x = Inches(0.4) + i * Inches(4.25)
        _rect(slide, x, Inches(1.15), Inches(4.05), Inches(5.5), CARD)
        _rect(slide, x, Inches(1.15), Inches(4.05), Inches(0.12), RED if i == 0 else NAVY2)
        _textbox(slide, x + Inches(0.2), Inches(1.4), Inches(3.65), Inches(0.8), item["title"], 16, True, NAVY)
        _textbox(slide, x + Inches(0.2), Inches(2.3), Inches(3.65), Inches(0.4), "Наблюдение", 11, True, RED)
        _textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3.65), Inches(1.2), item["fact"], 14, False, INK)
        _textbox(slide, x + Inches(0.2), Inches(3.95), Inches(3.65), Inches(0.35), "Что это значит", 11, True, RED)
        _textbox(slide, x + Inches(0.2), Inches(4.3), Inches(3.65), Inches(1.0), item["meaning"], 14, False, INK)
        _textbox(slide, x + Inches(0.2), Inches(5.4), Inches(3.65), Inches(0.35), "Что дал бы полный уровень", 11, True, RED)
        _textbox(slide, x + Inches(0.2), Inches(5.75), Inches(3.65), Inches(0.7), item["full"], 13, False, INK)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def _media_count(media: dict[str, int], catalog: Catalog, canonical: str) -> int:
    aliases = [canonical]
    for outlet in catalog.outlets:
        if outlet["canonical"] == canonical:
            aliases.extend(outlet.get("aliases") or [])
    n = 0
    for name, c in media.items():
        low = name.lower()
        if any(a.lower() in low or low in a.lower() for a in aliases):
            n += c
    return n


def _speaker_bars(bundle: ReportBundle, catalog: Catalog) -> tuple[list[str], list[int]]:
    from collections import Counter
    total: Counter[str] = Counter()
    for st in bundle.object_stats:
        total.update(st.speakers)
    labels, vals = [], []
    for pid, n in total.most_common(10):
        labels.append(catalog.person_by_id[pid].short if pid in catalog.person_by_id else pid)
        vals.append(n)
    return labels, vals


def _add_table(slide, rows: list[list[str]], l, t, w, h):
    if not rows:
        return
    table = slide.shapes.add_table(len(rows), len(rows[0]), l, t, w, h).table
    for j in range(len(rows[0])):
        table.columns[j].width = int(w / len(rows[0]))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if i else 11)
                p.font.bold = i == 0
                p.font.color.rgb = WHITE if i == 0 else INK
                p.font.name = "Calibri"
            fill = NAVY if i == 0 else (CARD if i % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
